#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""生成 EJIT 进展汇报 PPT —— 通用汇报风格(文件名按当前日期动态命名)。"""
from datetime import datetime
from pathlib import Path
import subprocess
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ───────────────── 通用汇报风格配色(蓝灰系·舒适)─────────────────
PRIMARY      = RGBColor(0x2C, 0x52, 0x82)   # 深蓝 主强调(标题/竖条/数据)
ACCENT       = RGBColor(0x31, 0x82, 0xCE)   # 钢蓝 次强调
DARK         = RGBColor(0x2D, 0x37, 0x48)   # 深蓝灰 表头
SLATE        = RGBColor(0x4A, 0x55, 0x68)   # 中深灰蓝 层名块
GRAY         = RGBColor(0x59, 0x59, 0x59)   # 中灰 正文
GRAY_L       = RGBColor(0x8C, 0x8C, 0x8C)   # 浅灰 副文本
BG_GRAY      = RGBColor(0xF5, 0xF7, 0xFA)   # 浅蓝灰背景
BORDER       = RGBColor(0xD9, 0xDE, 0xE6)   # 边框灰蓝
ROW_ALT      = RGBColor(0xF7, 0xF9, 0xFC)   # 表格行交替
PRIM_L       = RGBColor(0xEB, 0xF2, 0xFA)   # 浅蓝(仅★核心价值)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)

# 状态徽章色(柔和, 仅小色块使用)
GREEN        = RGBColor(0x38, 0xA1, 0x6F)
YELLOW       = RGBColor(0xD6, 0x9E, 0x2E)
ORANGE       = RGBColor(0xDD, 0x6B, 0x20)
STATUS_COLOR = {"✓": GREEN, "⚙": YELLOW, "△": ORANGE}

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
FONT = "微软雅黑"

REPO_ROOT = Path(__file__).resolve().parents[1]
BASE_REF_CANDIDATES = ("release/21.x", "origin/release/21.x")


def git(*args):
    return subprocess.check_output(["git", *args], cwd=REPO_ROOT, text=True)


def resolve_base_ref():
    for ref in BASE_REF_CANDIDATES:
        try:
            subprocess.check_output(
                ["git", "rev-parse", "--verify", ref],
                cwd=REPO_ROOT,
                text=True,
                stderr=subprocess.DEVNULL,
            )
            return ref
        except subprocess.CalledProcessError:
            pass
    return BASE_REF_CANDIDATES[0]


def _add_count(value):
    return 0 if value == "-" else int(value)


def collect_work_stats():
    """统计当前分支相对 release/21.x 的全量增量。"""
    base_ref = resolve_base_ref()
    range_spec = f"{base_ref}...HEAD"
    numstat = git(
        "diff", "--numstat", range_spec, "--", ".",
        ":(exclude)jit_design_doc/*.pptx",
    )
    stats = {
        "added": 0,
        "code_added": 0,
        "source_added": 0,
        "test_added": 0,
        "doc_added": 0,
        "design_docs": 0,
        "commits": 0,
    }
    source_exts = {
        ".c", ".cc", ".cpp", ".cxx", ".h", ".hpp", ".inc", ".td",
        ".def", ".cmake", ".py", ".sh", ".ld", ".ll",
    }
    for line in numstat.splitlines():
        added_s, _deleted_s, path = line.split("\t")[:3]
        added = _add_count(added_s)
        path = path.replace("\\", "/")
        suffix = Path(path).suffix
        is_test = path.startswith("ejit_test/") or "/test/" in path or "/unittests/" in path
        is_doc = path.startswith("jit_design_doc/") or path == "CLAUDE.md" or path.endswith("README.md")
        is_code = suffix in source_exts or path.endswith("CMakeLists.txt")
        stats["added"] += added
        if is_code or (is_test and not path.endswith("README.md")):
            stats["code_added"] += added
        if is_code and not is_doc:
            stats["source_added"] += added
        if is_test:
            stats["test_added"] += added
        if is_doc:
            stats["doc_added"] += added

    doc_files = git(
        "diff", "--name-only", range_spec, "--", "jit_design_doc",
        ":(exclude)jit_design_doc/*.pptx",
    )
    stats["design_docs"] = sum(1 for p in doc_files.splitlines() if p.endswith(".md"))
    stats["commits"] = int(git("rev-list", "--count", f"{base_ref}..HEAD").strip())
    return stats


WORK_STATS = collect_work_stats()


def fmt_k(n):
    return f"{n / 1000:.1f}K" if n >= 10000 else f"{n:,}"

def slide(): return prs.slides.add_slide(BLANK)

def rect(s, x, y, w, h, color, line=None, line_w=0.5):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is not None:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp

def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=None):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing: p.space_after = Pt(spacing)
        for (t, sz, col, b, f) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col
            r.font.bold = b; r.font.name = f
    return tb

def card(s, x, y, w, h, bar=PRIMARY):
    """统一白底 + 细灰边框 + 左侧主色竖条。"""
    rect(s, x, y, w, h, WHITE, line=BORDER, line_w=0.75)
    rect(s, x, y, Inches(0.07), h, bar)

def header(s, title, subtitle=None):
    rect(s, 0, 0, SW, Inches(0.07), PRIMARY)
    rect(s, 0, Inches(0.07), Inches(0.45), Inches(0.55), PRIMARY)
    txt(s, Inches(0.6), Inches(0.12), Inches(11.0), Inches(0.5),
        [[(title, 23, DARK, True, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        txt(s, Inches(0.6), Inches(0.58), Inches(12.4), Inches(0.26),
            [[(subtitle, 10.5, GRAY, False, FONT)]])
    rect(s, 0, Inches(0.88), SW, Inches(0.015), BORDER)

# ════════════════════════════════════════════════════════════════
# 第 1 页：整体架构完成情况
# ════════════════════════════════════════════════════════════════
def page1():
    s = slide()
    header(s, "EJIT 整体架构与完成情况",
           "EmbeddedJIT · 时间窗常量 + 运行时特化 · 2026 上半年")

    lx = Inches(0.35); ly = Inches(0.95); lw = Inches(8.5)

    # 范式条 — 深蓝底白字
    rect(s, lx, ly, lw, Inches(0.42), PRIMARY)
    txt(s, lx, ly, lw, Inches(0.42),
        [[("① 范式: 时间窗常量 + 运行时特化", 12, WHITE, True, FONT),
          ("    ✓", 12, WHITE, True, FONT)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, lx + lw/2 - Inches(0.1), ly + Inches(0.42), Inches(0.2), Inches(0.18),
        [[("▼", 10, GRAY_L, True, FONT)]], align=PP_ALIGN.CENTER)

    # 源码层 — 深灰底
    yy = ly + Inches(0.6)
    rect(s, lx, yy, lw, Inches(0.5), SLATE)
    txt(s, lx + Inches(0.15), yy, lw - Inches(0.3), Inches(0.5),
        [[("源码层  ", 12, WHITE, True, FONT),
          ("② Clang 6 属性标注面 ", 11.5, WHITE, True, FONT),
          ("(ejit_entry/_period/_period_arr/_period_arr_ind/_may_const/_period_lc)", 9.5, RGBColor(0xCC,0xCC,0xCC), False, FONT),
          ("  ✓", 12, WHITE, True, FONT)]],
        anchor=MSO_ANCHOR.MIDDLE)
    txt(s, lx + lw/2 - Inches(0.1), yy + Inches(0.5), Inches(0.2), Inches(0.18),
        [[("▼", 10, GRAY_L, True, FONT)]], align=PP_ALIGN.CENTER)

    # AOT 编译期 — 白底红竖条
    yy = ly + Inches(1.32)
    card(s, lx, yy, lw, Inches(0.98))
    txt(s, lx + Inches(0.18), yy + Inches(0.04), lw - Inches(0.3), Inches(0.3),
        [[("AOT 编译期 (两阶段)", 12, DARK, True, FONT)]])
    txt(s, lx + Inches(0.18), yy + Inches(0.34), lw - Inches(0.3), Inches(0.3),
        [[("O2/O3前:  ", 10, DARK, True, FONT),
          ("③ 两阶段bitcode提取", 10, GRAY, False, FONT),("✓",10,GREEN,True,FONT),
          ("   ④ may_const双重保留+预优化", 10, GRAY, False, FONT),("✓",10,GREEN,True,FONT),
          ("   ⑤ 协调器PASS5", 10, GRAY, False, FONT),("✓",10,GREEN,True,FONT)]])
    txt(s, lx + Inches(0.18), yy + Inches(0.61), lw - Inches(0.3), Inches(0.3),
        [[("O2/O3后:  ", 10, DARK, True, FONT),
          ("⑥ Period注册", 10, GRAY, False, FONT),("✓",10,GREEN,True,FONT),
          ("   ⑦ 单函数Wrapper+稠密索引", 10, GRAY, False, FONT),("✓",10,GREEN,True,FONT),
          ("   ⑧ 生命周期护栏", 10, GRAY, False, FONT),("✓",10,GREEN,True,FONT),
          ("   ⑨ 双轨注册", 10, GRAY, False, FONT),("✓",10,GREEN,True,FONT)]])
    txt(s, lx + lw/2 - Inches(0.1), yy + Inches(0.98), Inches(0.2), Inches(0.18),
        [[("▼ @__ejit_bitcode + 弱符号表", 9, GRAY_L, True, FONT)]], align=PP_ALIGN.CENTER)

    # 运行时特化引擎 — 白底红竖条
    yy = ly + Inches(2.52)
    card(s, lx, yy, lw, Inches(1.08))
    txt(s, lx + Inches(0.18), yy + Inches(0.04), lw - Inches(0.3), Inches(0.3),
        [[("运行时特化引擎", 12, DARK, True, FONT),
          ("  业务调用 → wrapper jit_entry → 命中缓存跳转 / 失败 fallback AOT", 9.5, GRAY_L, False, FONT)]])
    txt(s, lx + Inches(0.18), yy + Inches(0.35), lw - Inches(0.3), Inches(0.3),
        [[("⑩ 缓存LRU", 10, DARK, True, FONT),("✓",10,GREEN,True,FONT),
          (" → ⑪ 调度器", 10, DARK, True, FONT),("✓",10,GREEN,True,FONT),
          (" → ⑫ Bitcode加载器", 10, DARK, True, FONT),("✓",10,GREEN,True,FONT),
          (" → ⑬ 优化流水线", 10, DARK, True, FONT),("✓",10,GREEN,True,FONT)]])
    txt(s, lx + Inches(0.18), yy + Inches(0.63), lw - Inches(0.3), Inches(0.3),
        [[("→ ⑭ PASS6 结构体字段常量替换", 10, DARK, True, FONT),
          (" ★核心价值",10,PRIMARY,True,FONT),("✓",10,GREEN,True,FONT),
          ("  → ⑮ OrcJIT(LLJIT+JITLink)", 10, DARK, True, FONT),("✓",10,GREEN,True,FONT)]])
    txt(s, lx + Inches(0.18), yy + Inches(0.86), lw - Inches(0.3), Inches(0.3),
        [[("代码内存: ⑯ EJitCodePool (SRE 4K内存池, 可选启用)", 10, DARK, True, FONT),
          ("  ✓", 10, GREEN, True, FONT)]])
    txt(s, lx + lw/2 - Inches(0.1), yy + Inches(1.08), Inches(0.2), Inches(0.18),
        [[("▼", 10, GRAY_L, True, FONT)]], align=PP_ALIGN.CENTER)

    # 异步 Taskpool — 白底红竖条
    yy = ly + Inches(3.78)
    card(s, lx, yy, lw, Inches(0.78))
    txt(s, lx + Inches(0.18), yy + Inches(0.04), lw - Inches(0.3), Inches(0.3),
        [[("异步编译调度 / Taskpool", 12, DARK, True, FONT),
          ("  (当前分支重点)", 9.5, GRAY, False, FONT)]])
    txt(s, lx + Inches(0.18), yy + Inches(0.36), lw - Inches(0.3), Inches(0.3),
        [[("⑰ 单worker调度", 10, DARK, True, FONT),("✓",10,GREEN,True,FONT),
          ("   ⑱ 跨核共享单worker", 10, DARK, True, FONT),("△",10,ORANGE,True,FONT),
          ("   ⑲ 无锁队列+去重+分桶缓存", 10, DARK, True, FONT),("✓",10,GREEN,True,FONT),
          ("   ⑳ version失效", 10, DARK, True, FONT),("✓",10,GREEN,True,FONT)]])

    # 横切支柱 — 白底红竖条
    yy = ly + Inches(4.72)
    card(s, lx, yy, lw, Inches(0.98))
    txt(s, lx + Inches(0.18), yy + Inches(0.04), lw - Inches(0.3), Inches(0.3),
        [[("横切支柱 (贯穿所有层)", 12, DARK, True, FONT)]])
    txt(s, lx + Inches(0.18), yy + Inches(0.35), lw - Inches(0.3), Inches(0.3),
        [[("㉑ 稠密槽注册表", 10, DARK, True, FONT),("✓",10,GREEN,True,FONT),
          ("   ㉒ 注册暂存+冻结", 10, DARK, True, FONT),("✓",10,GREEN,True,FONT),
          ("   ㉓ SRE平台抽象", 10, DARK, True, FONT),("△",10,ORANGE,True,FONT),
          ("   ㉔ 并发原语层", 10, DARK, True, FONT),("✓",10,GREEN,True,FONT)]])
    txt(s, lx + Inches(0.18), yy + Inches(0.63), lw - Inches(0.3), Inches(0.3),
        [[("㉕ LLVM裁剪三层", 10, DARK, True, FONT),("✓",10,GREEN,True,FONT),
          ("  (53.a 117MB→单ejit.o 32MB)   ", 9, GRAY_L, False, FONT),
          ("㉖ BareMetal裸核", 10, DARK, True, FONT),("✓",10,GREEN,True,FONT),
          ("  (no-op头+POSIX桩)", 9, GRAY_L, False, FONT)]])

    # ── 右：完成情况统计 ──
    rx = Inches(9.0); ry = Inches(0.95); rw = Inches(4.0)
    rect(s, rx, ry, rw, Inches(0.42), SLATE)
    txt(s, rx, ry, rw, Inches(0.42),
        [[("完成情况统计", 13, WHITE, True, FONT)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 两个状态卡(去掉默认OFF和被取代)
    sy = ry + Inches(0.52)
    # ✓ 已实现
    card(s, rx, sy, rw, Inches(0.85), bar=GREEN)
    txt(s, rx + Inches(0.2), sy + Inches(0.06), Inches(2.0), Inches(0.55),
        [[("24", 30, DARK, True, FONT)]])
    txt(s, rx + Inches(1.5), sy + Inches(0.12), rw - Inches(1.7), Inches(0.35),
        [[("✓ 已实现使用", 12, GREEN, True, FONT)]])
    txt(s, rx + Inches(1.5), sy + Inches(0.46), rw - Inches(1.7), Inches(0.3),
        [[("代码完成并接入", 9.5, GRAY, False, FONT)]])
    # △ 待真机
    sy2 = sy + Inches(0.95)
    card(s, rx, sy2, rw, Inches(0.85), bar=ORANGE)
    txt(s, rx + Inches(0.2), sy2 + Inches(0.06), Inches(2.0), Inches(0.55),
        [[("2", 30, DARK, True, FONT)]])
    txt(s, rx + Inches(1.5), sy2 + Inches(0.12), rw - Inches(1.7), Inches(0.35),
        [[("△ 代码完成待真机", 12, ORANGE, True, FONT)]])
    txt(s, rx + Inches(1.5), sy2 + Inches(0.46), rw - Inches(1.7), Inches(0.3),
        [[("待 aarch64 SRE多核验证", 9.5, GRAY, False, FONT)]])

    # 总技术点红条
    sy3 = sy2 + Inches(0.95)
    rect(s, rx, sy3, rw, Inches(0.5), PRIMARY)
    txt(s, rx + Inches(0.15), sy3, rw - Inches(0.3), Inches(0.5),
        [[("26 个技术点  ·  覆盖 6 个域", 12, WHITE, True, FONT)]],
        anchor=MSO_ANCHOR.MIDDLE)

    # 工作量数据(含测试)
    sy4 = sy3 + Inches(0.6)
    txt(s, rx, sy4, rw, Inches(0.3),
        [[("▎工作量数据", 12, DARK, True, FONT)]])
    work = [
        (fmt_k(WORK_STATS["added"]), "全部增量"),
        (fmt_k(WORK_STATS["code_added"]), "代码增量"),
        (fmt_k(WORK_STATS["test_added"]), "测试/验证"),
        (str(WORK_STATS["commits"]), "分支提交"),
        (str(WORK_STATS["design_docs"]), "设计文档"),
        ("aarch64大端", "SRE裸核"),
        ("6+6", "Clang属性与Pass"),
        ("2", "编译模式"),
    ]
    wy = sy4 + Inches(0.35)
    for i, (v, k) in enumerate(work):
        col = i % 2; row = i // 2
        x = rx + col * Inches(2.05)
        y = wy + row * Inches(0.5)
        rect(s, x, y, Inches(1.95), Inches(0.45), BG_GRAY, line=BORDER, line_w=0.5)
        txt(s, x + Inches(0.12), y, Inches(1.05), Inches(0.45),
            [[(v, 13, PRIMARY, True, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
        txt(s, x + Inches(1.05), y, Inches(0.88), Inches(0.45),
            [[(k, 8.5, GRAY, False, FONT)]], anchor=MSO_ANCHOR.MIDDLE)

    # 底部时间轴(2026 上半年开发节奏)
    tl_y = Inches(6.7)
    # 节点均匀分布在 [1.3, 12.0],首末留足标签边距
    node_start = Inches(1.3); node_end = Inches(12.0)
    # 时间轴主线
    line_y = tl_y + Inches(0.66)
    rect(s, node_start, line_y, node_end - node_start, Inches(0.035), PRIMARY)
    # 9 个里程碑节点(标签在上、日期在下,均在节点上方)
    milestones = [
        ("05.09", "项目立项"),
        ("05.15", "设计文档首版"),
        ("05.20", "AOT pass主体"),
        ("05.25", "运行时特化"),
        ("06.10", "aarch64 bare-metal裁剪"),
        ("06.15", "SRE code pool"),
        ("06.16", "demo打通"),
        ("06.18", "4K内存池+taskpool"),
        ("06.27", "跨核共享worker"),
    ]
    n = len(milestones)
    seg = (node_end - node_start) / (n - 1)
    for i, (date, label) in enumerate(milestones):
        cx = node_start + seg * i
        # 标签(最上,单行)
        txt(s, cx - Inches(0.62), tl_y, Inches(1.24), Inches(0.3),
            [[(label, 8, DARK, True, FONT)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM)
        # 日期(节点正上方)
        txt(s, cx - Inches(0.4), tl_y + Inches(0.32), Inches(0.8), Inches(0.2),
            [[(date, 8.5, PRIMARY, True, FONT)]], align=PP_ALIGN.CENTER)
        # 节点圆点
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, cx - Inches(0.055), line_y - Inches(0.04), Inches(0.11), Inches(0.11))
        dot.fill.solid(); dot.fill.fore_color.rgb = PRIMARY
        dot.line.color.rgb = WHITE; dot.line.width = Pt(1.2)
        dot.shadow.inherit = False

# ════════════════════════════════════════════════════════════════
# 第 2 页：详细技术项
# ════════════════════════════════════════════════════════════════
def page2():
    s = slide()
    header(s, "EJIT 详细技术项", "26 个技术点 · 6 个域 · 实现状态 · 关键设计决策")

    data = [
        ("①", "时间窗常量+运行时特化", "全局", "✓", "立身之本: 时间窗内数据当编译期常量特化, 失败回退AOT"),
        ("②", "Clang 6 属性标注面", "前端", "✓", "may_const 软标注(丢失安全); Attr.td→SemaEJIT→CGEJIT 全链路"),
        ("③", "两阶段 bitcode 提取嵌入 PASS1", "AOT", "✓", "O2前提取调用图闭包, embed @__ejit_bitcode (提取改CloneModule+删)"),
        ("④", "may_const 双重保留 + AOT预优化", "AOT", "✓", "LLVM层kind保留 + GV偏移回退(v1.7) + reAnnotate; PASS1↔PASS6纽带"),
        ("⑤", "AOT 协调器 PASS5", "AOT", "✓", "调度PASS2→3→4 + hasAnyEjitMetadata快速返回 + 一致性诊断"),
        ("⑥", "Period/静态变量注册 PASS2", "AOT", "✓", "生成 ejit_register_period_array / static_var 注册调用"),
        ("⑦", "单函数 Wrapper + 稠密索引 PASS3", "AOT", "✓", "演进: 名字哈希→运行时稠密funcIndex(防碰撞); JIT失败fall原函数体"),
        ("⑧", "时间窗生命周期护栏 PASS4", "AOT", "✓", "ejit_period_lc 函数入口deactivate / 出口activate, 同序配对"),
        ("⑨", "双轨运行时注册体系", "AOT", "✓", "global_ctors构造器 + __ejit_registry_*[]静态表(裸核); 超出设计"),
        ("⑩", "编译缓存 LRU", "运行时", "✓", "iterator内嵌O(1); cacheKey=funcIdx<<32|dims; 三重上限+period失效"),
        ("⑪", "编译调度器 CompileDriver", "运行时", "✓", "热路径单哈希查缓存, miss走compileCold(解码/验证/编译/入缓存)"),
        ("⑫", "Bitcode 加载器 ModuleLoader", "运行时", "✓", "funcIdx稠密键(非funcName); 缓存period元数据; 幂等/冲突拒绝"),
        ("⑬", "JIT 优化流水线 EJitOptimizer", "运行时", "✓", "参数替换→InstCombine→PASS6→L1/L2/L3; Inline禁用(AOT预内联)"),
        ("⑭", "结构体字段常量替换 PASS6 ★", "运行时", "✓", "核心价值: may_const load→GEP算offset→读内存→构造Constant; 三模式"),
        ("⑮", "OrcJIT 引擎 EJitOrcEngine", "运行时", "✓", "LLJIT+JITLink替代MCJIT; per-cacheKey JITDylib隔离; Large code model"),
        ("⑯", "SRE 机器码内存池 EJitCodePool", "运行时", "✓", "2MiB池接管JITLink代码内存, 4K内存池(同页装512函数); callback注入可单测"),
        ("⑰", "纯异步单 worker 调度", "taskpool", "✓", "producer入队即返回fallback, 平台task上单worker轮询; CompileCallback解耦"),
        ("⑱", "跨核共享单 worker taskpool", "taskpool", "△", "POD共享blob+CAS选举owner; generation代际隔离; 代码✓待真机验证"),
        ("⑲", "无锁队列+去重+分桶缓存", "taskpool", "✓", "Vyukov MPSC(SC硬约束→固定1worker); 32桶隔离rehash; commit gate锁内校验"),
        ("⑳", "SwitchController version 失效", "taskpool", "✓", "逐实例version单调; toggle后三层全懒失效; 8×256二维数组零分配"),
        ("㉑", "稠密槽注册表管理", "横切", "✓", "注册期单线程分配稠密槽(lifecycle[0,8)/funcIndex[0,4096)), 无锁"),
        ("㉒", "注册暂存 + 双路径 + 冻结", "横切", "✓", "构造期暂存→init消费; 构造器/静态表双路径; taskpool下冻结注册"),
        ("㉓", "SRE 平台抽象层", "横切", "△", "头抽象+链接择一; 平台符号无weak fallback(缺失即链接错); host✓ SRE待接入"),
        ("㉔", "并发原语层", "横切", "✓", "EJitAtomic(__atomic_*)/EJitRwLock(双变量)/EJitIpcLock(短临界区)"),
        ("㉕", "LLVM 裁剪三层体系", "横切", "✓", "源码宏排除+EJitPassBuilder+lipo.py; 53.a(117MB)→单ejit.o 32MB"),
        ("㉖", "BareMetal 裸核运行时", "横切", "✓", "EJIT_FREESTANDING; no-op std头+POSIX桩库; std::mutex→BareMetalMutex"),
    ]

    cols = [("#", 0.42), ("技术点", 3.35), ("层", 0.85), ("状态", 0.62), ("关键设计决策", 7.05)]
    tx = Inches(0.3); ty = Inches(1.05)
    cx = tx
    for name, w in cols:
        rect(s, cx, ty, Inches(w), Inches(0.36), DARK)
        txt(s, cx, ty, Inches(w), Inches(0.36),
            [[(name, 11, WHITE, True, FONT)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += Inches(w)

    ry = ty + Inches(0.36)
    rh = Inches(0.197)
    for i, (no, name, layer, st, note) in enumerate(data):
        is_star = "★" in name
        bg = PRIM_L if is_star else (ROW_ALT if i % 2 == 0 else WHITE)
        cx = tx
        for name_, w in cols:
            rect(s, cx, ry, Inches(w), rh, bg, line=BORDER, line_w=0.25)
            cx += Inches(w)
        cx = tx
        txt(s, cx, ry, Inches(cols[0][1]), rh, [[(no, 9.5, PRIMARY, True, FONT)]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += Inches(cols[0][1])
        txt(s, cx + Inches(0.05), ry, Inches(cols[1][1]) - Inches(0.1), rh,
            [[(name, 9.3, DARK, True, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
        cx += Inches(cols[1][1])
        layer_color = {"全局":PRIMARY,"前端":DARK,"AOT":DARK,"运行时":DARK,
                       "taskpool":DARK,"横切":GRAY}.get(layer, GRAY)
        txt(s, cx, ry, Inches(cols[2][1]), rh, [[(layer, 8.5, layer_color, True, FONT)]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += Inches(cols[2][1])
        rect(s, cx + Inches(0.12), ry + Inches(0.025), Inches(0.34), rh - Inches(0.05), STATUS_COLOR[st])
        txt(s, cx + Inches(0.12), ry + Inches(0.025), Inches(0.34), rh - Inches(0.05),
            [[(st, 10, WHITE, True, FONT)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += Inches(cols[3][1])
        note_color = PRIMARY if is_star else GRAY
        txt(s, cx + Inches(0.05), ry, Inches(cols[4][1]) - Inches(0.1), rh,
            [[(note, 8.6, note_color, False, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
        ry += rh

    # 底部图例(去掉被取代说明)
    legy = ry + Inches(0.1)
    rect(s, Inches(0.3), legy, Inches(12.7), Inches(0.4), BG_GRAY, line=BORDER, line_w=0.5)
    txt(s, Inches(0.5), legy, Inches(12.5), Inches(0.4),
        [[("图例:  ", 10, DARK, True, FONT),
          ("✓", 11, GREEN, True, FONT),(" 已实现使用(24)   ", 9.5, GRAY, False, FONT),
          ("△", 11, ORANGE, True, FONT),(" 代码完成待真机验证(2)   ", 9.5, GRAY, False, FONT),
          ("    ★ 核心价值", 9.5, PRIMARY, True, FONT)]],
        anchor=MSO_ANCHOR.MIDDLE)

# ════════════════════════════════════════════════════════════════
# 第 3 页：EJIT 后续工作规划
# ════════════════════════════════════════════════════════════════
def page3_roadmap():
    s = slide()
    header(s, "EJIT 后续工作规划", "从代码完成 → 真机验证 → 性能收益 → 可交付工程化")

    bx = Inches(0.35); by = Inches(0.98); bw = Inches(12.65)
    rect(s, bx, by, bw, Inches(0.5), PRIMARY)
    txt(s, bx, by, bw, Inches(0.5),
        [[("目标: 形成可复现的 EJIT 闭环证据链 —— ", 11.5, WHITE, True, FONT),
          ("真机可跑", 11.5, RGBColor(0xCF,0xE0,0xF5), True, FONT),
          (" / ", 11.5, WHITE, True, FONT),
          ("收益可量化", 11.5, RGBColor(0xCF,0xE0,0xF5), True, FONT),
          (" / ", 11.5, WHITE, True, FONT),
          ("交付可裁剪", 11.5, RGBColor(0xCF,0xE0,0xF5), True, FONT)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ── 三阶段路线图 ──
    phases = [
        ("近期：验证闭环", "2~4 周", GREEN, [
            ("SRE/AArch64 真机", "验证跨核共享 worker、裸核链接、平台符号接入"),
            ("Taskpool 压测", "MPSC 队列、去重、分桶缓存、version 失效长稳测试"),
            ("JITLink/CodePool", "4K/2MiB 池边界、释放策略、失败 fallback 行为"),
            ("端到端 Demo", "从属性标注到运行时特化的可复现实验脚本"),
        ]),
        ("中期：收益量化", "1~2 月", ACCENT, [
            ("Benchmark 体系", "zlib/zstd/业务 kernel: cold/hot/命中率/编译耗时"),
            ("PASS6 增强", "数组/嵌套结构/volatile/别名场景覆盖与收益分析"),
            ("调参策略", "L1/L2/L3 pipeline、缓存上限、异步编译触发阈值"),
            ("裁剪体积", "单 ejit.o 继续瘦身: PassBuilder/Target/JITLink 依赖收敛"),
        ]),
        ("远期：工程交付", "季度", PRIMARY, [
            ("平台抽象固化", "host/SRE/bare-metal 三套配置稳定化与文档化"),
            ("诊断与可观测", "trace/dump/错误码/统计 counters 统一输出"),
            ("安全护栏", "注册冻结、符号白名单、内存权限切换与越界保护"),
            ("论文/专利材料", "时间窗常量 + 运行时特化 + 裁剪 JIT 的体系化总结"),
        ]),
    ]

    x0 = Inches(0.35); y0 = Inches(1.65); gap = Inches(0.2)
    colw = (Inches(12.65) - gap * 2) / 3
    for idx, (title, span, color, items) in enumerate(phases):
        x = x0 + idx * (colw + gap)
        rect(s, x, y0, colw, Inches(0.42), color)
        txt(s, x + Inches(0.12), y0, colw - Inches(0.24), Inches(0.42),
            [[(title, 12, WHITE, True, FONT), ("  ·  ", 10, WHITE, False, FONT),
              (span, 10, RGBColor(0xE8,0xF2,0xFF), True, FONT)]],
            anchor=MSO_ANCHOR.MIDDLE)
        cy = y0 + Inches(0.55)
        for no, (name, desc) in enumerate(items, 1):
            card(s, x, cy, colw, Inches(0.72), bar=color)
            rect(s, x + Inches(0.16), cy + Inches(0.14), Inches(0.34), Inches(0.34), color)
            txt(s, x + Inches(0.16), cy + Inches(0.14), Inches(0.34), Inches(0.34),
                [[(str(no), 9.5, WHITE, True, FONT)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            txt(s, x + Inches(0.58), cy + Inches(0.07), colw - Inches(0.75), Inches(0.26),
                [[(name, 10.5, DARK, True, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
            txt(s, x + Inches(0.58), cy + Inches(0.34), colw - Inches(0.75), Inches(0.32),
                [[(desc, 8.3, GRAY, False, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
            cy += Inches(0.8)

    # ── 底部交付物 ──
    dy = Inches(6.15)
    txt(s, Inches(0.35), dy - Inches(0.28), Inches(12.65), Inches(0.25),
        [[("▎阶段性交付物", 12.5, DARK, True, FONT)]])
    deliverables = [
        ("真机验证报告", "SRE 多核 / 裸核 / fallback"),
        ("性能收益表", "命中率 / 耗时 / 加速比"),
        ("裁剪构建包", "LLVMEJIT.a / 单 ejit.o"),
        ("可复现 Demo", "脚本 + 测试 + 文档"),
    ]
    dw = Inches(3.05)
    for i, (name, desc) in enumerate(deliverables):
        x = Inches(0.35) + i * Inches(3.2)
        rect(s, x, dy, dw, Inches(0.72), BG_GRAY, line=BORDER, line_w=0.5)
        rect(s, x, dy, Inches(0.06), Inches(0.72), PRIMARY)
        txt(s, x + Inches(0.18), dy + Inches(0.06), dw - Inches(0.3), Inches(0.28),
            [[(name, 10.5, DARK, True, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
        txt(s, x + Inches(0.18), dy + Inches(0.36), dw - Inches(0.3), Inches(0.28),
            [[(desc, 8.5, GRAY, False, FONT)]], anchor=MSO_ANCHOR.MIDDLE)

# ════════════════════════════════════════════════════════════════
# 第 4 页：编译器的技术规划与畅想 (AI × 传统编译器)
# ════════════════════════════════════════════════════════════════
def page4():
    s = slide()
    header(s, "编译器的技术规划与畅想", "AI × 传统编译器 · 用 AI 改造编译器本身")

    # 定位条
    bx = Inches(0.35); by = Inches(0.98); bw = Inches(12.65)
    rect(s, bx, by, bw, Inches(0.5), PRIMARY)
    txt(s, bx, by, bw, Inches(0.5),
        [[("用 AI 改造传统编译器本身（而非用编译器服务 AI）· 两大落点：", 11.5, WHITE, True, FONT),
          ("① 产出更优代码（性能）", 11.5, RGBColor(0xCF,0xE0,0xF5), True, FONT),
          ("   ② 自身/产物更小更快（小型化）", 11.5, RGBColor(0xCF,0xE0,0xF5), True, FONT)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ── 左栏: 渐进落地 9 方向 ──
    lx = Inches(0.35); ly = Inches(1.62); lw = Inches(6.25)
    txt(s, lx, ly, lw, Inches(0.3),
        [[("▎渐进落地方向 (9) · 可工程化、能发论文+工业界会用", 12.5, DARK, True, FONT)]])

    perf = [
        ("1", "学习型 Inliner 可解释化", "蒸馏 MLGO 黑盒 → 决策树规则, 可审计/可移植"),
        ("2", "LLM Superoptimizer ★", "LLM propose + Alive2 verify, 重写 10-30 条指令"),
        ("3", "神经代价模型", "替代 RegAlloc/Scheduler 启发式(异构硬件撑不住)"),
        ("4", "PGO 无 profile 化", "GNN 预测分支频次, 让 PGO 民主化"),
        ("5", "自动发掘 Combine 规则 ★", "扫手写汇编 diff, 自动产真实 LLVM patch"),
    ]
    size = [
        ("6", "学习型 Size Opt", "-Oz 继承者, MLGO 换 size reward"),
        ("7", "LTO 学习型死代码消除", "预测运行时不可达 symbol(冷代码后置)"),
        ("8", "AI 辅助 Pass 删除 / 瘦身", "定制 PassManager + 裁剪编译器 binary"),
        ("9", "IR 学习型压缩", "bitcode/MIR 专用字典, 分布式编译加速"),
    ]
    iy = ly + Inches(0.35)
    txt(s, lx, iy, lw, Inches(0.25),
        [[("性能优化（让编译器产出更优代码）", 10.5, ACCENT, True, FONT)]])
    iy += Inches(0.27)
    for no, name, desc in perf:
        card(s, lx, iy, lw, Inches(0.42), bar=ACCENT)
        txt(s, lx + Inches(0.18), iy, Inches(0.3), Inches(0.42),
            [[(no, 11, PRIMARY, True, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
        star = "★" in name
        txt(s, lx + Inches(0.5), iy, Inches(2.5), Inches(0.42),
            [[(name, 9.5, DARK, True, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
        txt(s, lx + Inches(2.95), iy, lw - Inches(3.05), Inches(0.42),
            [[(desc, 8.2, GRAY, False, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
        iy += Inches(0.45)
    iy += Inches(0.05)
    txt(s, lx, iy, lw, Inches(0.25),
        [[("小型化（编译器自身/产物更小更快）", 10.5, ACCENT, True, FONT)]])
    iy += Inches(0.27)
    for no, name, desc in size:
        card(s, lx, iy, lw, Inches(0.42), bar=ACCENT)
        txt(s, lx + Inches(0.18), iy, Inches(0.3), Inches(0.42),
            [[(no, 11, PRIMARY, True, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
        txt(s, lx + Inches(0.5), iy, Inches(2.5), Inches(0.42),
            [[(name, 9.5, DARK, True, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
        txt(s, lx + Inches(2.95), iy, lw - Inches(3.05), Inches(0.42),
            [[(desc, 8.2, GRAY, False, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
        iy += Inches(0.45)

    # ── 右栏: 颠覆性畅想精选 ──
    rx = Inches(6.8); ry = Inches(1.62); rw = Inches(6.2)
    txt(s, rx, ry, rw, Inches(0.3),
        [[("▎颠覆性畅想 (精选) · 触及编译器根基假设", 12.5, DARK, True, FONT)]])

    radical = [
        ("夯", "自演化编译器", "编译时持续学习优化规则并固化进下版本, 传统编译器 5 年内或被淘汰", PRIMARY),
        ("夯", "Verified-by-Default", "AI 自动生成性质证明, 代码默认带形式化证书(无 UB/data race)", PRIMARY),
        ("顶级", "端到端可微编译器", "整条 pipeline 离散决策连续松弛, 端到端梯度优化运行时间", ACCENT),
        ("顶级", "程序综合替代编译", "编译器重写算法本身(O(n²)→O(n log n)), I/O 等价+profile 验证", ACCENT),
        ("顶级", "神经符号混合 IR", "IR 节点可为神经嵌入, 中间阶段“模糊思考”后塌缩成符号", ACCENT),
        ("人上人", "跨层联合优化", "打破 编译器/OS/硬件 分层抽象, AI 学跨层隐性接口", GREEN),
        ("人上人", "多模态编译器", "读代码+注释+commit+issue, 理解“程序意图”作优化信号", GREEN),
    ]
    iy = ry + Inches(0.35)
    for tier, name, desc, c in radical:
        card(s, rx, iy, rw, Inches(0.62), bar=c)
        # 评级徽章
        rect(s, rx + Inches(0.18), iy + Inches(0.1), Inches(0.55), Inches(0.42), c)
        txt(s, rx + Inches(0.18), iy + Inches(0.1), Inches(0.55), Inches(0.42),
            [[(tier, 8, WHITE, True, FONT)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, rx + Inches(0.85), iy + Inches(0.05), rw - Inches(1.0), Inches(0.3),
            [[(name, 11, DARK, True, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
        txt(s, rx + Inches(0.85), iy + Inches(0.32), rw - Inches(1.0), Inches(0.28),
            [[(desc, 8.3, GRAY, False, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
        iy += Inches(0.66)

    # ── 底部: 推荐切入路径 + EJIT 结合 ──
    bby = Inches(6.62)
    rect(s, bx, bby, bw, Inches(0.78), BG_GRAY, line=BORDER, line_w=0.5)
    rect(s, bx, bby, Inches(0.07), Inches(0.78), PRIMARY)
    txt(s, bx + Inches(0.18), bby + Inches(0.04), bw - Inches(0.3), Inches(0.28),
        [[("▎推荐切入路径", 11, DARK, True, FONT)]])
    txt(s, bx + Inches(0.18), bby + Inches(0.3), bw - Inches(0.3), Inches(0.48),
        [[("① 快速出成果建声誉 → 自动挖 Combine 规则(每周产真实 LLVM patch)   ", 9, GRAY, False, FONT),
          ("② 立大功 → LLM Superoptimizer(范式干净、Alive2 兜底)   ", 9, GRAY, False, FONT),
          ("③ 长跑 → 自演化 / Verified-by-Default", 9, GRAY, False, FONT)]],
        anchor=MSO_ANCHOR.MIDDLE)

page1()
page2()
page3_roadmap()
page4()

_today = datetime.now().strftime("%Y%m%d")
out = f"/home/ruanchen/djq/github/ejit/llvm-project/jit_design_doc/EJIT进展_{_today}.pptx"
prs.save(out)
print("已生成:", out)
print("共 4 页 (通用汇报风格)")
